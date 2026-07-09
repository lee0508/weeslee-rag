# -*- coding: utf-8 -*-
"""
Semantic structure extraction and semantic chunk preparation.

현재는 PPTX 문서의 의미 섹션 구조화를 우선 지원한다.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from app.services.document_structure_extractor import get_document_structure_extractor
from app.services.rfp_pattern_analyzer import rfp_pattern_analyzer


TOP_SECTION_RULES = [
    ("표지", ["표지", "추진계획", "용역 추진계획"]),
    ("목차", ["목차", "목차 내용"]),
    ("사업개요", ["사업개요", "추진배경 및 목적", "추진배경", "목적", "사업범위"]),
    ("현황 및 문제점", ["현황 및 문제점", "서비스 현황", "정보화 현황", "선진사례", "문제점 및 개선 방향"]),
    ("세부내용", ["세부내용", "추진현황", "협력사업 추진현황"]),
    ("주요내용", ["주요내용", "조사내용", "개선방안", "활용실태"]),
    ("사업 추진방안", ["사업 추진방안", "추진 목표", "추진 폭표", "추진전략", "추진체계", "추진일정"]),
    ("제안요청 내용", ["제안요청 내용", "제안요청 개요", "용어의 정의", "상세 요구사항", "요구사항 총괄", "요구사항 상세", "컨설팅 요구사항", "프로젝트 관리 요구사항", "프로젝트 지원 요구사항", "제약사항 요구사항", "보안 요구사항"]),
    ("사업자 선정방법", ["사업자 선정방법", "입찰 방법", "입찰 참가자격", "제안서 평가 및 낙찰자 결정", "제안서 평가방법"]),
    ("제안서 작성안내", ["제안서 작성안내", "제안서 작성 요령", "제안서의 효력", "제안서 효력", "제안서 작성 지침", "제안서 작성지침", "제안서 목차", "제안서 작성방법", "제안서 세부 작성 지침", "제안 안내 사항", "제안서 제출서류", "기타 사항", "기타사항"]),
    ("안내 및 유의사항", ["안내 및 유의사항", "제안요청 설명회", "사업자의 의무 및 책임에 관한 사항", "지식재산 공동귀속", "SW산출물 반출 절차", "지식재산권 공동 활용 범위에 관한 사항", "기타 유의사항"]),
    ("붙임 및 별지", ["붙임", "별첨", "별표", "별지서식", "입찰 관련 서식", "기술평가항목 및 배점표", "입찰평가항목 및 배점표", "외주 용역사업 보안 특약 조항", "소프트웨어사업 영향평가 검토결과서"]),
    ("추진일정 및 체계", ["추진일정 및 체계", "사업 추진일정", "추진체계", "사업 추진보고", "주요 산출물"]),
    ("사업이해도", ["사업이해도", "사업 이해", "사업환경", "사업 환경", "배경", "목적", "범위", "전제조건", "특장점", "특징", "장점", "목표시스템", "구성도", "구성체계"]),
    ("추진전략", ["추진전략", "핵심성공", "핵심 질의", "핵심질의", "추진조직", "추진체계", "협력방안", "품질 확보", "전문협업", "전략 1", "전략 2", "전략 3", "전략 4"]),
    ("컨설팅 방법론", ["방법론", "WIM2", "ISP 방법론", "ISMP 방법론", "프로젝트 준비", "환경분석", "현황분석", "목표모델", "이행계획", "실행계획", "산출물", "적용 사례", "절차 및 활동내용"]),
]

SUBSECTION_ALIASES = {
    "표지 내용": ["표지 내용", "표지", "추진계획", "용역 추진계획"],
    "목차": ["목차", "목차 내용"],
    "사업개요": ["사업개요"],
    "추진배경 및 목적": ["추진배경 및 목적", "추진배경", "목적"],
    "사업범위": ["사업범위", "사업 범위"],
    "기대효과 및 성과지표": ["기대효과 및 성과지표", "기대효과", "성과지표", "주요 전산화 방향 및 기대효과"],
    "현황 및 문제점": ["현황 및 문제점"],
    "서비스 현황": ["서비스 현황"],
    "정보화 현황": ["정보화 현황"],
    "선진사례": ["선진사례"],
    "문제점 및 개선 방향": ["문제점 및 개선 방향", "개선 방향"],
    "세부내용": ["세부내용", "세부 내용"],
    "추진현황": ["추진현황", "협력사업 추진현황", "AFSIS 협력사업 추진현황", "농정원 AFSIS 협력사업 추진현황"],
    "주요내용": ["주요내용", "주요 내용", "활용실태", "개선방안 도출"],
    "사업 추진방안": ["사업 추진방안"],
    "추진 목표": ["추진 목표", "추진목표", "추진 폭표"],
    "추진전략": ["추진전략", "사업 추진전략"],
    "추진일정 및 체계": ["추진일정 및 체계", "추진 일정 및 체계"],
    "사업 추진일정": ["사업 추진일정", "추진일정"],
    "추진체계": ["추진체계", "사업 추진체계"],
    "사업 추진보고": ["사업 추진보고", "추진보고"],
    "주요 산출물": ["주요 산출물", "산출물"],
    "제안요청 내용": ["제안요청 내용", "제안요청내용"],
    "제안요청 개요": ["제안요청 개요"],
    "용어의 정의": ["용어의 정의"],
    "상세 요구사항": ["상세 요구사항", "요구사항 상세"],
    "요구사항 총괄": ["요구사항 총괄", "요구사항 구분 및 총괄표"],
    "컨설팅 요구사항": ["컨설팅 요구사항"],
    "프로젝트 관리 요구사항": ["프로젝트 관리 요구사항"],
    "프로젝트 지원 요구사항": ["프로젝트 지원 요구사항"],
    "제약사항 요구사항": ["제약사항 요구사항"],
    "보안 요구사항": ["보안 요구사항"],
    "사업자 선정방법": ["사업자 선정방법"],
    "입찰 방법": ["입찰 방법"],
    "입찰 참가자격": ["입찰 참가자격"],
    "제안서 평가 및 낙찰자 결정": ["제안서 평가 및 낙찰자 결정"],
    "제안서 평가방법": ["제안서 평가방법"],
    "제안서 작성안내": ["제안서 작성안내", "제안서 작성 요령"],
    "제안서의 효력": ["제안서의 효력", "제안서 효력"],
    "제안서 작성지침": ["제안서 작성지침", "제안서 작성 지침"],
    "제안서 목차": ["제안서 목차"],
    "제안서 작성방법": ["제안서 작성방법"],
    "제안서 세부 작성 지침": ["제안서 세부 작성 지침"],
    "제안 안내 사항": ["제안 안내 사항"],
    "제안서 제출서류": ["제안서 제출서류"],
    "기타 사항": ["기타 사항", "기타사항"],
    "안내 및 유의사항": ["안내 및 유의사항"],
    "제안요청 설명회": ["제안요청 설명회"],
    "사업자의 의무 및 책임에 관한 사항": ["사업자의 의무 및 책임에 관한 사항"],
    "지식재산 공동귀속": ["지식재산 공동귀속"],
    "SW산출물 반출 절차 등": ["SW산출물 반출 절차 등", "SW산출물 반출 절차"],
    "지식재산권 공동 활용 범위에 관한 사항": ["지식재산권 공동 활용 범위에 관한 사항"],
    "기타 유의사항": ["기타 유의사항"],
    "붙임 및 별지": ["붙임", "별첨", "별표", "별지서식"],
    "사업 환경의 이해": ["사업환경의 이해", "사업 환경의 이해", "사업환경의이해"],
    "사업의 배경 및 목적": ["사업의 배경 및 목적", "사업 배경 및 목적", "사업배경 및 목적", "배경 및 필요성", "사업목적", "추진목적", "추진방향"],
    "제안의 범위": ["제안의 범위", "사업수행범위", "제안 범위"],
    "전제조건": ["전제조건", "제안의 전제조건"],
    "제안의 특징 및 장점": ["제안의 특징 및 장점", "제안의 특장점", "특장점"],
    "목표시스템 구성도": ["목표시스템 구성도", "목표 구성도"],
    "목표시스템 구성체계": ["목표시스템 구성체계", "목표 구성체계"],
    "핵심 질의기반 사업 추진전략 도출": ["핵심 질의기반 사업 추진전략 도출", "핵심질의기반 사업 추진전략 도출", "핵심질의", "핵심성공요소 기반의 4대 사업추진전략"],
    "사업 추진전략": ["사업 추진전략", "사업추진전략", "전략 1", "전략 2", "전략 3", "전략 4"],
    "추진체계": ["추진체계", "추진조직", "사업수행 조직", "업무 분장", "협력방안", "역할 분담"],
    "품질 확보방안": ["품질 확보방안", "최종산출물의 품질 확보방안"],
    "방법론 특성": ["방법론 특성", "컨설팅 방법론", "사업추진 방법론", "WIM2 BPR/ISP 방법론", "WIM2방법론의 특장점", "WIM2 방법론의 특장점"],
    "제안사 방법론 주요 적용 사례": ["제안사 방법론 주요 적용 사례", "방법론을 적용한 수행 프로젝트", "유사 수행 프로젝트"],
    "WIM2 방법론 절차 개요": ["WIM2 방법론 절차", "ISP 방법론 절차", "ISMP 방법론 절차", "방법론 절차 개요"],
    "프로젝트 준비": ["프로젝트 준비", "프로젝트 준비절차 및 활동내용"],
    "환경분석": ["환경분석", "환경분석 절차 및 활동내용"],
    "현황분석": ["현황분석", "현황분석 절차 및 활동내용"],
    "목표모델 수립": ["목표모델 수립", "목표모델수립", "목표모델 수립 절차 및 활동내용"],
    "이행계획 수립": ["이행계획 수립", "이행계획수립", "통합 실행계획 수립", "통합 실행계획 수립절차 및 활동내용", "실행계획수립"],
    "방법론 절차에 따른 산출물 예시": ["방법론 절차에 따른 산출물 예시", "단계별 산출물 내역", "환경분석 산출물", "현황분석 산출물", "목표모델수립 산출물", "이행계획수립 산출물"],
    "산출물 제출": ["산출물 제출", "최종산출물 제출", "단계별 활동내역 및 산출물 정보"],
}

STOPWORDS = {
    "전략 및 방법론",
    "전략및 방법론",
    "전략및방법론",
    "예 시",
    "추후보완",
    "활용 도구",
    "및 기법",
}


def _normalize_space(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _clean_line(text: str) -> str:
    text = _normalize_space(text)
    text = text.replace("\u3000", " ")
    return text.strip(" -\u2022\t")


def _normalize_for_compare(text: str) -> str:
    return re.sub(r"[\s\-\u2022:()<>]", "", text).lower()


def _first_line(text: str) -> str:
    for line in text.split("\n"):
        cleaned = _clean_line(line)
        if cleaned:
            return cleaned
    return ""


def _extract_sequence_index(text: str) -> Optional[Tuple[int, int]]:
    match = re.search(r"\((\d+)\s*/\s*(\d+)\)", text)
    if not match:
        return None
    return int(match.group(1)), int(match.group(2))


def _get_shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return _clean_line(shape.text)


def _extract_slide_title(slide) -> str:
    if getattr(slide.shapes, "title", None) is not None:
        title_text = _get_shape_text(slide.shapes.title)
        if title_text:
            return _first_line(title_text)

    candidates: List[Tuple[int, int, str]] = []
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if not text:
            continue
        top = getattr(shape, "top", 10**9)
        left = getattr(shape, "left", 10**9)
        line = _first_line(text)
        if line:
            candidates.append((top, left, line))

    if not candidates:
        return "제목 미확인"

    candidates.sort(key=lambda item: (item[0], item[1], len(item[2])))
    return candidates[0][2]


def _extract_slide_items(slide, title: str) -> List[str]:
    items: List[str] = []
    seen = set()
    title_norm = _normalize_for_compare(title)

    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if not text:
            continue
        lines = [_clean_line(line) for line in text.split("\n")]
        lines = [line for line in lines if line]
        for line in lines:
            line_norm = _normalize_for_compare(line)
            if not line_norm or line_norm == title_norm or line_norm in seen:
                continue
            seen.add(line_norm)
            items.append(line)

        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                row_values = [_clean_line(cell.text) for cell in row.cells if _clean_line(cell.text)]
                if row_values:
                    line = " | ".join(row_values)
                    line_norm = _normalize_for_compare(line)
                    if line_norm and line_norm not in seen:
                        seen.add(line_norm)
                        items.append(line)

    return items


def _infer_top_section(title: str, section_name: str) -> str:
    title_norm = _normalize_for_compare(title)
    for label, keywords in TOP_SECTION_RULES:
        for keyword in keywords:
            if _normalize_for_compare(keyword) in title_norm:
                return label
    return section_name or "기타"


def _canonicalize_subsection(text: str) -> Optional[str]:
    cleaned = _clean_line(text)
    if not cleaned:
        return None
    cleaned = re.sub(r"\(\d+/\d+\)", "", cleaned).strip()
    cleaned = re.sub(r"^[0-9IVXⅠⅡⅢⅣⅤ]+\.*\s*", "", cleaned).strip()
    cleaned = re.sub(r"^[0-9]+\.[0-9]+\s*", "", cleaned).strip()
    cleaned_norm = _normalize_for_compare(cleaned)

    for canonical, aliases in SUBSECTION_ALIASES.items():
        for alias in [canonical] + aliases:
            if _normalize_for_compare(alias) in cleaned_norm:
                return canonical
    return None


def _extract_subsection_label(slide: Dict[str, Any]) -> str:
    candidates = [str(item) for item in slide["items"][:8]] + [str(slide["title"])]
    for candidate in candidates:
        label = _canonicalize_subsection(candidate)
        if label:
            return label
    return str(slide["title"])


def _extract_detail_label(slide: Dict[str, Any], subsection: str) -> str:
    candidates = [str(item) for item in slide["items"][:30]] + [str(slide["title"])]
    sequence = None
    for source in [str(item) for item in slide["items"][:5]] + [str(slide["title"])]:
        sequence = _extract_sequence_index(source)
        if sequence:
            break

    if subsection == "방법론 절차에 따른 산출물 예시":
        if sequence and sequence[1] == 4:
            labels = {
                1: "환경분석 산출물",
                2: "현황분석 산출물",
                3: "목표모델수립 산출물",
                4: "이행계획수립 산출물",
            }
            if sequence[0] in labels:
                return labels[sequence[0]]

    if subsection == "WIM2 방법론 절차 개요":
        if sequence and sequence[1] == 6:
            labels = {
                1: "WIM2 방법론 절차 개요",
                2: "프로젝트 준비",
                3: "환경분석",
                4: "현황분석",
                5: "목표모델 수립",
                6: "통합 실행계획 수립",
            }
            if sequence[0] in labels:
                return labels[sequence[0]]

    for candidate in candidates:
        cleaned = _clean_line(candidate)
        cleaned = re.sub(r"\(\d+/\d+\)", "", cleaned).strip()
        cleaned = re.sub(r"^[0-9IVXⅠⅡⅢⅣⅤ]+\.*\s*", "", cleaned).strip()
        cleaned = re.sub(r"^[0-9]+\.[0-9]+\s*", "", cleaned).strip()
        if cleaned and _normalize_for_compare(cleaned) != _normalize_for_compare(subsection):
            return cleaned
    return subsection


def _slide_number_list(slide_numbers: List[int]) -> List[int]:
    return sorted(set(int(num) for num in slide_numbers))


def _summarize_slide_numbers(slide_numbers: List[int]) -> str:
    if not slide_numbers:
        return "슬라이드 미상"
    numbers = _slide_number_list(slide_numbers)
    ranges: List[str] = []
    start = numbers[0]
    end = numbers[0]
    for num in numbers[1:]:
        if num == end + 1:
            end = num
            continue
        ranges.append(str(start) if start == end else f"{start}~{end}")
        start = end = num
    ranges.append(str(start) if start == end else f"{start}~{end}")
    return "슬라이드 " + ", ".join(ranges)


def _collect_keywords(texts: List[str], limit: int = 12) -> List[str]:
    seen = set()
    keywords: List[str] = []
    for text in texts:
        for raw_part in re.split(r"[,\|/]| - |:|;|\n", text):
            part = _clean_line(raw_part)
            if not part or part in STOPWORDS:
                continue
            part = re.sub(r"^\d+(\.\d+)*\s*", "", part).strip()
            part = re.sub(r"\(\d+/\d+\)", "", part).strip()
            if not part or len(part) < 2 or len(part) > 40 or part in STOPWORDS:
                continue
            key = _normalize_for_compare(part)
            if not key or key in seen:
                continue
            seen.add(key)
            keywords.append(part)
            if len(keywords) >= limit:
                return keywords
    return keywords


def _resolve_section_group(relative_path: str, fallback: str = "기타") -> str:
    rel_obj = Path(relative_path) if relative_path else None
    if rel_obj and len(rel_obj.parts) >= 3:
        return rel_obj.parts[2]
    if rel_obj and len(rel_obj.parts) >= 2:
        return rel_obj.parts[1]
    return fallback


def _extract_page_content_items(page_text: str, page_title: str, limit: int = 8) -> List[str]:
    title_norm = _normalize_for_compare(page_title)
    items: List[str] = []
    for raw_line in str(page_text or "").splitlines():
        cleaned = _clean_line(raw_line)
        if not cleaned:
            continue
        if _normalize_for_compare(cleaned) == title_norm:
            continue
        if len(cleaned) < 2:
            continue
        items.append(cleaned)
        if len(items) >= limit:
            break
    return items


def _extract_slides(file_path: str) -> List[Dict[str, Any]]:
    prs = Presentation(file_path)
    slides: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _extract_slide_title(slide)
        items = _extract_slide_items(slide, title)
        slides.append({"slide_no": idx, "title": title, "items": items})
    return slides


def _build_outline(slides: List[Dict[str, Any]], section_name: str) -> List[Dict[str, Any]]:
    groups: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    for slide in slides:
        top_section = _infer_top_section(str(slide["title"]), section_name)
        if current is None or current["top_section"] != top_section:
            current = {"top_section": top_section, "slides": []}
            groups.append(current)
        current["slides"].append(slide)
    return groups


def _merge_group_slides(group_slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    merged: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    for slide in group_slides:
        subsection = _extract_subsection_label(slide)
        if current is None or current["subsection"] != subsection:
            current = {
                "subsection": subsection,
                "slide_numbers": [],
                "items": [],
                "detail_entries": [],
                "_seen": set(),
                "_detail_seen": set(),
            }
            merged.append(current)

        current["slide_numbers"].append(int(slide["slide_no"]))
        detail_label = _extract_detail_label(slide, subsection)
        detail_key = (detail_label, int(slide["slide_no"]))
        if detail_key not in current["_detail_seen"]:
            current["_detail_seen"].add(detail_key)
            current["detail_entries"].append({"label": detail_label, "slide_no": int(slide["slide_no"])})
        for item in slide["items"]:
            norm = _normalize_for_compare(str(item))
            if not norm or norm in current["_seen"]:
                continue
            current["_seen"].add(norm)
            current["items"].append(item)

    for item in merged:
        item.pop("_seen", None)
        item.pop("_detail_seen", None)
    return merged


def build_pptx_structure(file_path: str, relative_path: str = "") -> Dict[str, Any]:
    path_obj = Path(file_path)
    rel_path = relative_path or path_obj.name
    rel_obj = Path(rel_path)
    section_name = rel_obj.parts[2] if len(rel_obj.parts) >= 3 else (rel_obj.parts[1] if len(rel_obj.parts) >= 2 else "기타")
    slides = _extract_slides(file_path)
    groups = _build_outline(slides, section_name)
    sections: List[Dict[str, Any]] = []

    for idx, group in enumerate(groups, start=1):
        top_section = str(group["top_section"])
        merged_slides = _merge_group_slides(group["slides"])
        top_slide_numbers: List[int] = []
        subsection_entries: List[Dict[str, Any]] = []
        top_keyword_texts = [top_section]

        for sub_idx, slide in enumerate(merged_slides, start=1):
            subsection = str(slide["subsection"])
            slide_numbers = _slide_number_list(list(slide["slide_numbers"]))
            items = [str(item) for item in slide["items"]]
            detail_entries = slide.get("detail_entries", [])
            top_slide_numbers.extend(slide_numbers)
            detail_titles = [str(entry["label"]) for entry in detail_entries]
            top_keyword_texts.extend([subsection] + items[:10] + detail_titles)
            subsection_entries.append({
                "section_id": f"{idx}.{sub_idx}",
                "section_name": subsection,
                "parent_section": f"{idx}. {top_section}",
                "slide_range": [slide_numbers[0], slide_numbers[-1]],
                "slide_numbers": slide_numbers,
                "slide_label": _summarize_slide_numbers(slide_numbers),
                "subsections": [{"title": title, "slide_numbers": [int(entry["slide_no"])]} for title, entry in zip(detail_titles, detail_entries)],
                "content_items": items,
                "keywords": _collect_keywords([subsection] + items + detail_titles),
            })

        top_slide_numbers = _slide_number_list(top_slide_numbers)
        sections.append({
            "section_id": str(idx),
            "section_name": top_section,
            "slide_range": [top_slide_numbers[0], top_slide_numbers[-1]] if top_slide_numbers else [],
            "slide_numbers": top_slide_numbers,
            "slide_label": _summarize_slide_numbers(top_slide_numbers),
            "keywords": _collect_keywords(top_keyword_texts),
            "subsections": subsection_entries,
        })

    return {
        "file_name": path_obj.name,
        "source_path": str(path_obj),
        "relative_path": rel_obj.as_posix(),
        "section_group": section_name,
        "document_type": "pptx",
        "total_slides": len(slides),
        "structure_mode": "semantic_sections",
        "sections": sections,
    }


def build_text_semantic_structure(
    text: str,
    *,
    document_id: Optional[int] = None,
    file_name: str = "",
    relative_path: str = "",
    file_type: str = "",
) -> Dict[str, Any]:
    """
    OCR/파싱 결과 텍스트에서 범용 semantic structure를 생성한다.

    - cover/toc/page 구조를 우선 보존
    - Step 5의 chunk_semantic_sections()가 바로 사용할 수 있는 형태로 맞춘다
    """
    normalized_text = _normalize_space(str(text or ""))
    rel_path = relative_path or file_name
    section_group = _resolve_section_group(rel_path)
    analysis = rfp_pattern_analyzer.analyze_text_content(normalized_text) if normalized_text else {}

    extractor = get_document_structure_extractor()
    doc_id = int(document_id or 0)
    doc_structure = extractor.extract_structure(
        text=normalized_text,
        document_id=doc_id,
        folder_path=rel_path,
        file_type=file_type.lower().lstrip(".") if file_type else "",
    )

    top_labels = {
        "cover": "표지",
        "toc": "목차",
        "content": section_group or "본문",
        "appendix": "부록",
    }
    grouped_pages: Dict[str, List[Dict[str, Any]]] = {}
    section_seq = 1

    toc_titles = analysis.get("toc", {}).get("section_titles", []) or []
    toc_leaf_titles = [{"title": title, "slide_numbers": []} for title in toc_titles[:12]]

    for page in doc_structure.pages:
        page_type = str(page.page_type or "content")
        top_section = top_labels.get(page_type, section_group or "본문")
        page_title = _clean_line(str(page.page_title or "")) or f"페이지 {page.page_no}"
        content_items = _extract_page_content_items(page.text_content, page_title)
        leaf_sections = toc_leaf_titles if page_type == "toc" and toc_leaf_titles else []
        keywords = _collect_keywords([page_title] + content_items + [item.get("title", "") for item in leaf_sections])

        grouped_pages.setdefault(top_section, []).append({
            "section_id": str(section_seq),
            "section_name": page_title,
            "page_no": page.page_no,
            "page_type": page_type,
            "content_items": content_items,
            "keywords": keywords,
            "leaf_sections": leaf_sections,
        })
        section_seq += 1

    sections: List[Dict[str, Any]] = []
    for top_index, (top_section, entries) in enumerate(grouped_pages.items(), start=1):
        top_page_numbers = [int(entry["page_no"]) for entry in entries]
        top_keyword_texts = [top_section]
        subsection_entries: List[Dict[str, Any]] = []

        for sub_index, entry in enumerate(entries, start=1):
            top_keyword_texts.extend([entry["section_name"]] + entry["content_items"] + [leaf.get("title", "") for leaf in entry["leaf_sections"]])
            subsection_entries.append({
                "section_id": f"{top_index}.{sub_index}",
                "section_name": entry["section_name"],
                "parent_section": f"{top_index}. {top_section}",
                "slide_range": [entry["page_no"], entry["page_no"]],
                "slide_numbers": [entry["page_no"]],
                "slide_label": f"페이지 {entry['page_no']}",
                "subsections": entry["leaf_sections"],
                "content_items": entry["content_items"],
                "keywords": entry["keywords"],
                "page_type": entry["page_type"],
            })

        sections.append({
            "section_id": str(top_index),
            "section_name": top_section,
            "slide_range": [min(top_page_numbers), max(top_page_numbers)] if top_page_numbers else [],
            "slide_numbers": sorted(set(top_page_numbers)),
            "slide_label": f"페이지 {min(top_page_numbers)}~{max(top_page_numbers)}" if len(set(top_page_numbers)) > 1 else (f"페이지 {top_page_numbers[0]}" if top_page_numbers else "페이지 미상"),
            "keywords": _collect_keywords(top_keyword_texts),
            "subsections": subsection_entries,
        })

    structured = {
        "file_name": file_name or Path(rel_path).name,
        "source_path": "",
        "relative_path": Path(rel_path).as_posix() if rel_path else "",
        "section_group": section_group,
        "document_type": file_type.lower().lstrip(".") if file_type else "",
        "total_slides": doc_structure.total_pages,
        "structure_mode": "semantic_sections",
        "sections": sections,
        "cover_page": analysis.get("cover_page", {}),
        "toc": analysis.get("toc", {}),
        "detected_sections": analysis.get("detected_sections", []),
        "document_summary": analysis.get("summary", ""),
        "page_types": [str(page.page_type or "content") for page in doc_structure.pages],
    }
    structured["semantic_tags"] = infer_semantic_tags(structured)
    return structured


def infer_semantic_tags(structured: Dict[str, Any]) -> Dict[str, Any]:
    joined = " ".join(structured.get("relative_path", "").split("/"))
    text_pool: List[str] = [joined]
    for top in structured.get("sections", []):
        text_pool.append(str(top.get("section_name") or ""))
        text_pool.extend(str(value) for value in top.get("keywords") or [])
        for sub in top.get("subsections", []):
            text_pool.append(str(sub.get("section_name") or ""))
            text_pool.extend(str(value) for value in sub.get("keywords") or [])
            for leaf in sub.get("subsections", []):
                text_pool.append(str(leaf.get("title") or ""))

    text = " ".join(text_pool)
    methodology = "WIM2" if "wim2" in text.lower() else ""
    domain = "감사" if "감사" in text else ""
    technology = "AI" if "AI" in text or "인공지능" in text else ""
    return {
        "methodology": methodology,
        "domain": domain,
        "technology": technology,
    }
