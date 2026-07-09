# PPTX 문서 추출기 - 그룹 재귀 및 이미지 OCR 보충
# 작업일: 2026-07-08 - DB 시스템 설정 연동 추가, structured_txt 우선 참조 추가
"""
PPTX extractor with group recursion and image OCR supplementation.
[2026-07-08] structured_txt 우선 사용 지원 - 수동 추출 파일이 있으면 먼저 사용
"""
import io
import logging
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.extractors.base import BaseExtractor, ExtractionResult

logger = logging.getLogger(__name__)

# [2026-07-08] structured 파일 우선 사용 지원
try:
    from app.services.structured_content_resolver import StructuredContentResolver
    HAS_STRUCTURED_RESOLVER = True
except ImportError:
    HAS_STRUCTURED_RESOLVER = False
    StructuredContentResolver = None


def _get_db_ocr_setting(key: str, default):
    """DB에서 OCR 설정 조회. 실패 시 default 반환."""
    try:
        from app.services.system_settings_service import get_system_setting
        return get_system_setting("ocr", key, default)
    except Exception:
        return default

try:
    from app.services.text_quality_checker import text_quality_checker
    HAS_QUALITY_CHECKER = True
except ImportError:
    HAS_QUALITY_CHECKER = False
    text_quality_checker = None


def _libreoffice_path() -> str:
    """Return the local soffice path."""
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if os.path.exists(candidate):
            return candidate
    return "soffice"


class PptxExtractor(BaseExtractor):
    """PPTX text extractor with group recursion and image OCR."""

    def _try_structured_txt(self, file_path: str) -> Optional[Dict[str, Any]]:
        """
        [2026-07-08] structured_txt 우선 사용.
        1) 같은 폴더에 동일 이름의 .txt 파일이 있으면 사용
        2) 없으면 StructuredContentResolver로 structured_txt 폴더 확인
        """
        file_path_obj = Path(file_path)
        file_name = file_path_obj.name

        # 1단계: 같은 폴더에 .txt 파일 확인
        txt_path = file_path_obj.with_suffix(".txt")
        if txt_path.exists():
            try:
                # UTF-8 BOM 또는 UTF-8로 읽기
                try:
                    content = txt_path.read_text(encoding="utf-8-sig")
                except UnicodeDecodeError:
                    content = txt_path.read_text(encoding="utf-8")

                if content and len(content.strip()) > 100:
                    logger.info(f"[PPTX] 같은 폴더 txt 사용: {txt_path.name}")
                    return ExtractionResult(
                        success=True,
                        content=content,
                        metadata={
                            "source": file_path,
                            "filename": file_name,
                            "used_paths": [str(txt_path)],
                        },
                        method="same_folder_txt"
                    ).to_dict()
            except Exception as e:
                logger.debug(f"[PPTX] 같은 폴더 txt 읽기 실패: {e}")

        # 2단계: StructuredContentResolver로 structured_txt 폴더 확인
        if not HAS_STRUCTURED_RESOLVER or not StructuredContentResolver:
            return None
        try:
            path_str = str(file_path).replace("\\", "/")

            # relative_path 추론
            relative_path = file_name
            for marker in ["00. RAG 소스/", "01. RFP/", "02. 제안서/", "03. 산출물/"]:
                if marker in path_str:
                    idx = path_str.find(marker)
                    relative_path = path_str[idx:]
                    break

            class _FakeDoc:
                def __init__(self, rp, fn):
                    self.relative_path = rp
                    self.file_name = fn

            resolver = StructuredContentResolver({
                "use_structured_txt": True,
                "use_structured_json": True,
                "prefer_structured_content": True,
                "max_text_chars": 50000,
            })

            content = resolver.resolve_document_content(_FakeDoc(relative_path, file_name))
            combined_text = content.get("combined_text") or ""

            if combined_text and len(combined_text.strip()) > 100:
                logger.info(f"[PPTX] structured_txt 우선 사용: {file_name}")
                return ExtractionResult(
                    success=True,
                    content=combined_text,
                    metadata={
                        "source": file_path,
                        "filename": file_name,
                        "used_paths": content.get("used_paths", []),
                    },
                    method="structured_txt_priority"
                ).to_dict()
            return None
        except Exception as e:
            logger.debug(f"[PPTX] structured_txt 조회 실패: {e}")
            return None

    def __init__(
        self,
        use_ocr_supplement: bool = True,
        quality_threshold: float = 0.5,
        ocr_dpi: int = None,
        ocr_language: str = None,
        ocr_engine: str = None,
        ocr_use_gpu: Optional[bool] = None,
        ocr_image_min_bytes: int = 3000,
    ):
        self.use_ocr_supplement = use_ocr_supplement
        self.quality_threshold = quality_threshold
        # DB 설정 우선, 없으면 하드코딩 fallback
        self.ocr_dpi = max(72, int(ocr_dpi or _get_db_ocr_setting("ocr_dpi", 300)))
        self.ocr_language = str(ocr_language or _get_db_ocr_setting("ocr_language", "kor+eng"))
        self.ocr_engine = str(ocr_engine or _get_db_ocr_setting("ocr_engine", "tesseract")).lower()
        self.ocr_use_gpu = ocr_use_gpu
        self.ocr_image_min_bytes = int(ocr_image_min_bytes or 3000)

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def _iter_all_shapes(self, shapes) -> Iterable[Any]:
        """Recursively flatten group shapes and yield leaf shapes."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_all_shapes(shape.shapes)
            else:
                yield shape

    async def extract(self, file_path: str) -> Dict[str, Any]:
        import time
        start_time = time.time()

        if not os.path.exists(file_path):
            return ExtractionResult(
                success=False,
                error=f"File not found: {file_path}",
            ).to_dict()

        # [2026-07-08] structured_txt 우선 사용 - 수동 추출 파일이 있으면 먼저 사용
        structured_result = self._try_structured_txt(file_path)
        if structured_result:
            structured_result["metadata"] = structured_result.get("metadata", {})
            structured_result["metadata"]["processing_time_sec"] = round(time.time() - start_time, 2)
            return structured_result

        metadata: Dict[str, Any] = {
            "source": file_path,
            "filename": Path(file_path).name,
            "extraction_attempts": [],
        }

        direct_result = self._extract_with_pptx(file_path)
        picture_count = int((direct_result.metadata or {}).get("picture_count") or 0)
        metadata["extraction_attempts"].append({
            "method": "python-pptx(group-recursive)",
            "success": direct_result.success,
            "text_length": len(direct_result.content) if direct_result.content else 0,
            "picture_count": picture_count,
        })

        quality = self._check_quality(direct_result.content or "")
        metadata["quality"] = quality
        quality_low = quality.get("quality_score", 0) < self.quality_threshold

        ocr_by_slide: Dict[int, str] = {}
        if self.use_ocr_supplement and (picture_count > 0 or quality_low):
            ocr_by_slide = self._ocr_embedded_images(file_path)
            metadata["extraction_attempts"].append({
                "method": "embedded-image-ocr",
                "success": bool(ocr_by_slide),
                "text_length": sum(len(value) for value in ocr_by_slide.values()),
                "ocr_slides": sorted(ocr_by_slide.keys()),
            })

        if ocr_by_slide:
            merged = self._merge_native_and_ocr(direct_result.content or "", ocr_by_slide)
            direct_result.metadata.update(metadata)
            direct_result.metadata["final_method"] = "pptx_group_recursive+image_ocr"
            direct_result.metadata["ocr_supplemented"] = True
            direct_result.metadata["ocr_engine"] = self._embedded_image_ocr_engine()
            direct_result.metadata["content_length"] = len(merged)
            return ExtractionResult(
                success=True,
                content=merged,
                metadata=direct_result.metadata,
                method="pptx_group_recursive+image_ocr",
            ).to_dict()

        if direct_result.success and direct_result.content:
            if quality_low and self.use_ocr_supplement:
                pdf_ocr_result = await self._extract_with_pdf_ocr(file_path)
                metadata["extraction_attempts"].append({
                    "method": "pdf_ocr_supplement",
                    "success": pdf_ocr_result.success if pdf_ocr_result else False,
                    "text_length": len(pdf_ocr_result.content) if pdf_ocr_result and pdf_ocr_result.content else 0,
                })
                if pdf_ocr_result and pdf_ocr_result.success and pdf_ocr_result.content:
                    merged_content = self._merge_content(direct_result.content, pdf_ocr_result.content)
                    merged_quality = self._check_quality(merged_content)
                    metadata["merged_quality"] = merged_quality
                    if merged_quality.get("quality_score", 0) > quality.get("quality_score", 0):
                        return ExtractionResult(
                            success=True,
                            content=merged_content,
                            metadata={
                                **direct_result.metadata,
                                **metadata,
                                "final_method": "pptx_group_recursive+pdf_ocr",
                                "ocr_supplemented": True,
                                "ocr_engine": (pdf_ocr_result.metadata or {}).get("ocr_engine") or self.ocr_engine,
                            },
                            method="pptx_group_recursive+pdf_ocr",
                        ).to_dict()

            direct_result.metadata.update(metadata)
            direct_result.metadata["final_method"] = "pptx_group_recursive"
            return direct_result.to_dict()

        pdf_ocr_result = await self._extract_with_pdf_ocr(file_path)
        if pdf_ocr_result and pdf_ocr_result.success and pdf_ocr_result.content:
            pdf_ocr_result.metadata.update(metadata)
            pdf_ocr_result.metadata["final_method"] = "pptx_pdf_ocr_only"
            pdf_ocr_result.metadata["ocr_required"] = True
            return pdf_ocr_result.to_dict()

        return ExtractionResult(
            success=False,
            error="모든 추출 방법 실패",
            metadata=metadata,
            method="pptx_all_failed",
        ).to_dict()

    def _extract_with_pptx(self, file_path: str) -> ExtractionResult:
        try:
            prs = Presentation(file_path)
            content_parts: List[str] = []
            slide_count = 0
            picture_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_text = [f"--- Slide {slide_num} ---"]

                for shape in self._iter_all_shapes(slide.shapes):
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            if any(row_text):
                                slide_text.append(" | ".join(row_text))
                        continue

                    if getattr(shape, "has_chart", False):
                        chart_lines = self._extract_chart_text(shape.chart)
                        if chart_lines:
                            slide_text.append("[차트] " + " / ".join(chart_lines))
                        continue

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_count += 1
                        continue

                    if getattr(shape, "has_text_frame", False):
                        text_value = shape.text_frame.text
                        if text_value and text_value.strip():
                            slide_text.append(text_value)

                content_parts.append("\n".join(slide_text))

            content = "\n\n".join(content_parts)
            return ExtractionResult(
                success=True,
                content=content,
                metadata={
                    "source": file_path,
                    "filename": Path(file_path).name,
                    "slides": slide_count,
                    "picture_count": picture_count,
                    "content_length": len(content),
                },
                method="python-pptx",
            )
        except Exception as e:
            return ExtractionResult(success=False, error=str(e), method="python-pptx")

    def _extract_chart_text(self, chart) -> List[str]:
        lines: List[str] = []
        try:
            if chart.has_title and chart.chart_title.text_frame.text.strip():
                lines.append(chart.chart_title.text_frame.text.strip())
        except Exception:
            pass
        try:
            categories = [str(category) for category in chart.plots[0].categories if str(category).strip()]
            if categories:
                lines.append("범주: " + ", ".join(categories))
        except Exception:
            pass
        try:
            for series in chart.series:
                if series.name and str(series.name).strip():
                    lines.append("계열: " + str(series.name))
        except Exception:
            pass
        return lines

    def _ocr_embedded_images(self, file_path: str) -> Dict[int, str]:
        results: Dict[int, str] = {}
        try:
            prs = Presentation(file_path)
        except Exception:
            return results

        for slide_num, slide in enumerate(prs.slides, 1):
            texts: List[str] = []
            for shape in self._iter_all_shapes(slide.shapes):
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    blob = shape.image.blob
                except Exception:
                    continue
                if len(blob) < self.ocr_image_min_bytes:
                    continue
                ocr_text = self._ocr_image_bytes(blob)
                if ocr_text:
                    texts.append(ocr_text)
            if texts:
                results[slide_num] = "\n".join(texts)

        return results

    def _ocr_image_bytes(self, blob: bytes) -> str:
        try:
            from PIL import Image
        except ImportError:
            return ""

        try:
            img = Image.open(io.BytesIO(blob))
            img.load()
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")
        except Exception:
            return ""

        if self.ocr_engine in {"easyocr", "olmocr"}:
            try:
                import numpy as np
                from app.extractors.pdf_extractor import _get_easyocr_reader

                reader = _get_easyocr_reader(self.ocr_use_gpu)
                results = reader.readtext(np.array(img))
                text = "\n".join(item[1] for item in results).strip()
                if text:
                    return text
            except Exception:
                pass

        try:
            import pytesseract
            return pytesseract.image_to_string(img, lang=self.ocr_language).strip()
        except Exception:
            return ""

    async def _extract_with_pdf_ocr(self, file_path: str) -> Optional[ExtractionResult]:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                cmd = [
                    _libreoffice_path(),
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    file_path,
                ]
                process = subprocess.run(cmd, capture_output=True, timeout=120)
                if process.returncode != 0:
                    return None

                pdf_files = list(Path(tmpdir).glob("*.pdf"))
                if not pdf_files:
                    return None

                from app.extractors.pdf_extractor import PDFExtractor

                pdf_result = await PDFExtractor(
                    use_ocr=True,
                    ocr_threshold=1,
                    ocr_use_gpu=self.ocr_use_gpu,
                    ocr_dpi=self.ocr_dpi,
                    ocr_language=self.ocr_language,
                    extract_tables=False,
                    table_as_markdown=False,
                    ocr_engine=self.ocr_engine,
                ).extract(str(pdf_files[0]))

                if not pdf_result.get("success") or not str(pdf_result.get("content") or "").strip():
                    return None

                metadata = dict(pdf_result.get("metadata") or {})
                metadata.update({
                    "source": file_path,
                    "filename": Path(file_path).name,
                    "content_length": len(str(pdf_result.get("content") or "")),
                    "slides": metadata.get("pages") or metadata.get("slides"),
                    "ocr_lang": self.ocr_language,
                    "ocr_dpi": self.ocr_dpi,
                    "ocr_engine": metadata.get("ocr_engine") or metadata.get("selected_method") or self.ocr_engine,
                    "pdf_converted_for_ocr": True,
                })

                return ExtractionResult(
                    success=True,
                    content=str(pdf_result.get("content") or ""),
                    metadata=metadata,
                    method=f"pptx_{pdf_result.get('method') or 'ocr'}",
                )
        except Exception as e:
            print(f"[WARN] OCR extraction failed for {file_path}: {e}")
            return None

    def _merge_native_and_ocr(self, native: str, ocr_by_slide: Dict[int, str]) -> str:
        native_slides = self._split_slides(native)
        merged_parts: List[str] = []
        all_slide_nums = sorted(set(native_slides.keys()) | set(ocr_by_slide.keys()))

        for slide_num in all_slide_nums:
            native_text = native_slides.get(slide_num, "")
            block = native_text if native_text else f"--- Slide {slide_num} ---"
            ocr_text = ocr_by_slide.get(slide_num, "")
            if ocr_text:
                native_norm = self._normalize(native_text)
                new_lines = []
                for line in ocr_text.splitlines():
                    line_stripped = line.strip()
                    if len(line_stripped) < 2:
                        continue
                    if self._normalize(line_stripped) in native_norm:
                        continue
                    new_lines.append(line_stripped)
                if new_lines:
                    block += "\n[이미지 OCR]\n" + "\n".join(new_lines)
            merged_parts.append(block)

        return "\n\n".join(merged_parts)

    def _merge_content(self, direct_text: str, ocr_text: str) -> str:
        def parse_slides(text: str) -> Dict[int, str]:
            slides: Dict[int, str] = {}
            current_slide = 0
            current_text: List[str] = []

            for line in text.split("\n"):
                if line.startswith("--- Slide "):
                    if current_slide > 0:
                        slides[current_slide] = "\n".join(current_text)
                    try:
                        current_slide = int(line.split("--- Slide ")[1].split(" ")[0])
                    except (ValueError, IndexError):
                        current_slide += 1
                    current_text = []
                else:
                    current_text.append(line)

            if current_slide > 0:
                slides[current_slide] = "\n".join(current_text)
            return slides

        direct_slides = parse_slides(direct_text)
        ocr_slides = parse_slides(ocr_text)
        all_slides = sorted(set(direct_slides.keys()) | set(ocr_slides.keys()))

        merged_parts = []
        for slide_num in all_slides:
            direct_content = direct_slides.get(slide_num, "")
            ocr_content = ocr_slides.get(slide_num, "")
            direct_clean = direct_content.replace("\n", "").replace(" ", "")
            ocr_clean = ocr_content.replace("\n", "").replace(" ", "")

            if len(ocr_clean) > len(direct_clean) * 1.2:
                selected = ocr_content
            else:
                selected = direct_content

            merged_parts.append(f"--- Slide {slide_num} ---\n{selected}")

        return "\n\n".join(merged_parts)

    @staticmethod
    def _split_slides(text: str) -> Dict[int, str]:
        slides: Dict[int, str] = {}
        current = 0
        buffer: List[str] = []

        for line in text.split("\n"):
            if line.startswith("--- Slide "):
                if current > 0:
                    slides[current] = "\n".join(buffer)
                try:
                    current = int(line.split("--- Slide ")[1].split(" ")[0])
                except (ValueError, IndexError):
                    current += 1
                buffer = [line]
            else:
                buffer.append(line)

        if current > 0:
            slides[current] = "\n".join(buffer)
        return slides

    @staticmethod
    def _normalize(text: str) -> str:
        return text.replace(" ", "").replace("\n", "").replace("\t", "")

    def _embedded_image_ocr_engine(self) -> str:
        if self.ocr_engine in {"easyocr", "olmocr"}:
            return "easyocr"
        return "tesseract"

    def _check_quality(self, text: str) -> dict:
        if not HAS_QUALITY_CHECKER or not text_quality_checker:
            text_len = len(text)
            korean_count = sum(1 for c in text if "가" <= c <= "힣")
            denominator = len(text.replace(" ", "")) or 1
            korean_ratio = korean_count / denominator
            return {
                "text_length": text_len,
                "korean_ratio": round(korean_ratio, 4),
                "quality_score": 0.7 if text_len > 100 and korean_ratio > 0.1 else 0.4,
                "decision": "use_direct_text" if text_len > 100 else "need_fallback",
            }

        result = text_quality_checker.check(text)
        return result.to_dict()


pptx_extractor = PptxExtractor()
